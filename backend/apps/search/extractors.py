"""
Estrazione testo da file per indicizzazione (FASE 12, FASE 37).
"""
import csv
import email
import os


def extract_text(file_path, mime_type=None):
    """
    Estrae testo da file. Ritorna stringa vuota se tipo non supportato.
    """
    if not file_path or not os.path.isfile(file_path):
        return ""
    mime_type = (mime_type or "").lower()
    lower_path = file_path.lower()
    try:
        if (mime_type.startswith("text/") and "csv" not in mime_type) or (
            not mime_type and lower_path.endswith(".txt")
        ):
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        if "pdf" in mime_type or (not mime_type and lower_path.endswith(".pdf")):
            from pypdf import PdfReader

            reader = PdfReader(file_path)
            parts = []
            for page in reader.pages:
                try:
                    parts.append(page.extract_text() or "")
                except Exception:
                    parts.append("")
            return "\n".join(parts)
        if "wordprocessingml" in mime_type or "msword" in mime_type or (
            not mime_type and (lower_path.endswith(".docx") or lower_path.endswith(".doc"))
        ):
            try:
                from docx import Document as DocxDocument

                doc = DocxDocument(file_path)
                return "\n".join(p.text for p in doc.paragraphs)
            except Exception:
                return ""
        if "spreadsheetml" in mime_type or lower_path.endswith(".xlsx"):
            from openpyxl import load_workbook

            wb = load_workbook(file_path, read_only=True, data_only=True)
            parts = []
            try:
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        row_text = " ".join(str(cell) for cell in row if cell is not None)
                        if row_text.strip():
                            parts.append(row_text)
            finally:
                wb.close()
            return "\n".join(parts)
        if "csv" in mime_type or lower_path.endswith(".csv"):
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                parts = []
                for i, row in enumerate(reader):
                    if i >= 5000:
                        break
                    row_text = " ".join(row)
                    if row_text.strip():
                        parts.append(row_text)
                return "\n".join(parts)
        if "presentation" in mime_type or lower_path.endswith(".pptx"):
            try:
                from pptx import Presentation

                prs = Presentation(file_path)
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            parts.append(shape.text)
                return "\n".join(parts)
            except ImportError:
                return ""
        if "message" in mime_type or lower_path.endswith(".eml"):
            with open(file_path, "rb") as f:
                msg = email.message_from_binary_file(f)
            parts = [msg.get("subject", ""), msg.get("from", ""), msg.get("to", "")]
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        parts.append(payload.decode("utf-8", errors="replace"))
            return "\n".join(p for p in parts if p)
        return ""
    except Exception:
        return ""
